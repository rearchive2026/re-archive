import os
import json
import shutil
from datetime import datetime
import fitz  # PyMuPDF
from pptx import Presentation
import google.generativeai as genai
from dotenv import load_dotenv

# Load environment variables
load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)
    model = genai.GenerativeModel('gemini-1.5-flash')
else:
    print("Warning: GEMINI_API_KEY not found. AI features will be disabled.")
    model = None

def extract_text_from_pdf(file_path):
    text = ""
    try:
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text()
    except Exception as e:
        print(f"Error extracting PDF {file_path}: {e}")
    return text

def extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"Error extracting PPTX {file_path}: {e}")
    return text

def get_ai_summary(text, category):
    if not model:
        return {
            "title": "Untitled Document",
            "summary": "AI summary unavailable.",
            "tags": [category],
            "content": text[:500] + "..."
        }
    
    prompt = f"""
    당신은 재건축 데이터 전문가입니다. 다음 텍스트는 '{category}' 카테고리의 문서 내용입니다.
    이 내용을 분석하여 다음 JSON 형식으로 응답해 주세요. 반드시 순수 JSON만 응답하세요.
    
    {{
        "title": "문서의 핵심 제목 (간결하게)",
        "summary": "문서의 핵심 내용을 3-4줄로 요약",
        "tags": ["태그1", "태그2"],
        "content": "Jekyll 마크다운 형식으로 보기 좋게 정리된 본문 내용 (표나 리스트 활용 권장)"
    }}
    
    텍스트 내용:
    {text[:10000]}  # 텍스트가 너무 길면 자름
    """
    
    try:
        response = model.generate_content(prompt)
        # Remove markdown code block if present
        cleaned_response = response.text.strip().replace("```json", "").replace("```", "")
        return json.loads(cleaned_response)
    except Exception as e:
        print(f"AI Error: {e}")
        return {
            "title": "AI Processing Failed",
            "summary": "Error occurred during AI analysis.",
            "tags": [category, "error"],
            "content": text[:500] + "..."
        }

def process_data():
    raw_data_root = 'raw_data'
    # 로컬 작업 시 바로 옆의 re-archive 폴더로 결과물을 보냅니다.
    output_root = '../re-archive'
    
    # Ensure directories exist
    dirs_to_make = [
        os.path.join(output_root, '_data'),
        os.path.join(output_root, '_posts'),
        os.path.join(output_root, '_pages'),
        os.path.join(output_root, 'assets', 'raw_docs')
    ]
    for d in dirs_to_make:
        os.makedirs(d, exist_ok=True)

    # Process each category folder
    categories = ['Notices', 'FAQ', 'Guides', 'URLs']
    for cat in categories:
        cat_path = os.path.join(raw_data_root, cat)
        if not os.path.exists(cat_path):
            continue
            
        for filename in os.listdir(cat_path):
            file_path = os.path.join(cat_path, filename)
            if os.path.isdir(file_path):
                continue
                
            print(f"Processing: {file_path}")
            
            # Extract text based on file type
            content_text = ""
            if filename.endswith('.pdf'):
                content_text = extract_text_from_pdf(file_path)
            elif filename.endswith('.pptx'):
                content_text = extract_text_from_pptx(file_path)
            elif filename.endswith(('.txt', '.url', '.md')):
                with open(file_path, 'r', encoding='utf-8') as f:
                    content_text = f.read()
            elif filename.endswith('.json'):
                # Handle existing JSON logic (masking)
                with open(file_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                processed_data = {k: v for k, v in data.items() if not k.startswith('_')}
                with open(os.path.join(output_root, '_data', filename), 'w', encoding='utf-8') as f:
                    json.dump(processed_data, f, indent=2, ensure_ascii=False)
                continue
            else:
                continue

            # Get AI Analysis
            analysis = get_ai_summary(content_text, cat)
            
            # Copy original file to assets
            asset_path = os.path.join('assets', 'raw_docs', filename)
            shutil.copy2(file_path, os.path.join(output_root, asset_path))
            
            # Generate Jekyll Post/Page
            date_str = datetime.now().strftime("%Y-%m-%d")
            clean_filename = "".join([c if c.isalnum() else "-" for c in filename.split('.')[0]])
            post_filename = f"{date_str}-{clean_filename}.md"
            
            with open(os.path.join(output_root, '_posts', post_filename), 'w', encoding='utf-8') as f:
                f.write('---\n')
                f.write(f'layout: single\n')
                f.write(f'title: "{analysis["title"]}"\n')
                f.write(f'date: {date_str}\n')
                f.write(f'categories: {cat}\n')
                f.write(f'tags: {json.dumps(analysis["tags"], ensure_ascii=False)}\n')
                f.write(f'last_modified_at: {datetime.now().isoformat()}\n')
                f.write('---\n\n')
                
                f.write(f'### 요약\n{analysis["summary"]}\n\n')
                f.write(f'---\n\n')
                f.write(f'{analysis["content"]}\n\n')
                f.write(f'---\n\n')
                f.write(f'**원본 파일 다운로드:** [ {filename} ]({{{{ site.baseurl }}}}/{asset_path})\n')

    # Basic root file copy
    for root_file in ['_config.yml', 'Gemfile', 'index.md']:
        if os.path.exists(root_file):
            shutil.copy2(root_file, os.path.join(output_root, root_file))
    
    # Copy _pages and assets (excluding raw_docs handled above)
    for folder in ['_pages', 'assets']:
        if os.path.exists(folder):
            dest = os.path.join(output_root, folder)
            if folder == 'assets':
                # Skip copy if dest already exists to avoid issues, or copy content
                if not os.path.exists(dest):
                    shutil.copytree(folder, dest)
                else:
                    # Sync logic or skip
                    pass
            else:
                if os.path.exists(dest):
                    shutil.rmtree(dest)
                shutil.copytree(folder, dest)

if __name__ == '__main__':
    process_data()
